#!/usr/bin/env python3
"""
Standalone PDF conversion core for the PDF Process plugin.

Converts a PDF into Word / Excel / PowerPoint. Designed to be packaged into a
single Windows executable (via PyInstaller) so end users need no Python runtime.

Usage:
    convert <word|excel|ppt> <input.pdf> <output_path>

Conversion strategy:
    word   -> pdf2docx (layout-preserving PDF -> DOCX)
    excel  -> pdf2docx table extraction -> openpyxl (falls back to page text)
    ppt    -> PyMuPDF renders each page to an image -> python-pptx (one per slide)
"""
import io
import os
import sys

SUPPORTED = ("word", "excel", "ppt")

# 1 point == 12700 EMU (English Metric Units used by Office Open XML).
EMU_PER_POINT = 12700
# Render DPI for PDF -> PPT page images.
PPT_RENDER_DPI = 150


def convert_word(input_pdf, output_path):
    """PDF -> DOCX using pdf2docx (preserves text, images and tables)."""
    from pdf2docx import Converter

    cv = Converter(input_pdf)
    try:
        cv.convert(output_path)
    finally:
        cv.close()
    return output_path


def _extract_tables(input_pdf):
    """Returns a list of tables (each a list of rows) detected by pdf2docx."""
    from pdf2docx import Converter

    cv = Converter(input_pdf)
    try:
        return cv.extract_tables() or []
    finally:
        cv.close()


def convert_excel(input_pdf, output_path):
    """PDF -> XLSX. Uses detected tables; falls back to per-page text lines."""
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    tables = _extract_tables(input_pdf)
    if tables:
        for idx, table in enumerate(tables, start=1):
            ws = wb.create_sheet(title=("Table %d" % idx)[:31])
            for row in table:
                ws.append(["" if cell is None else str(cell) for cell in row])
    else:
        import fitz

        doc = fitz.open(input_pdf)
        try:
            for i, page in enumerate(doc, start=1):
                ws = wb.create_sheet(title=("Page %d" % i)[:31])
                for line in page.get_text().splitlines():
                    ws.append([line])
        finally:
            doc.close()

    if not wb.sheetnames:
        wb.create_sheet(title="Sheet1")
    wb.save(output_path)
    return output_path


def convert_ppt(input_pdf, output_path, dpi=PPT_RENDER_DPI):
    """PDF -> PPTX. Each page is rendered to an image placed on its own slide."""
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(input_pdf)
    try:
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)

        for page_index, page in enumerate(doc):
            if page_index == 0:
                prs.slide_width = Emu(int(page.rect.width * EMU_PER_POINT))
                prs.slide_height = Emu(int(page.rect.height * EMU_PER_POINT))
            pixmap = page.get_pixmap(matrix=matrix)
            image = io.BytesIO(pixmap.tobytes("png"))
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image, 0, 0, width=prs.slide_width, height=prs.slide_height
            )

        prs.save(output_path)
    finally:
        doc.close()
    return output_path


_DISPATCH = {
    "word": convert_word,
    "excel": convert_excel,
    "ppt": convert_ppt,
}


def convert(fmt, input_pdf, output_path):
    """Dispatch a conversion by format name.

    Raises ValueError for an unsupported format and FileNotFoundError when the
    input PDF does not exist.
    """
    if fmt not in SUPPORTED:
        raise ValueError("Unsupported format: %s (expected one of %s)" % (fmt, ", ".join(SUPPORTED)))
    if not os.path.isfile(input_pdf):
        raise FileNotFoundError(input_pdf)
    out_dir = os.path.dirname(os.path.abspath(output_path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    return _DISPATCH[fmt](input_pdf, output_path)


def main(argv=None):
    """CLI entry point. Returns a process exit code (0 == success)."""
    argv = list(sys.argv[1:] if argv is None else argv)
    if len(argv) != 3:
        sys.stderr.write("usage: convert <word|excel|ppt> <input.pdf> <output_path>\n")
        return 2
    fmt, input_pdf, output_path = argv
    try:
        convert(fmt, input_pdf, output_path)
    except Exception as exc:  # noqa: BLE001 - surface any failure to the caller
        sys.stderr.write("ERROR: %s\n" % exc)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
